from tkinter import *
from tkinter import filedialog
from pptx import Presentation
import glob
from natsort import natsorted, ns

#define button commands

#select source folder
def select_1Click():
    global sourcePath
    photofolder = filedialog.askdirectory(initialdir = "S:\\Active Files\\PHOTOS\\STEWARDSHIP PHOTOS")
    sourcePath = photofolder
    inputPath.insert(0, sourcePath)

#select output folder
def select_2Click():
    global destPath
    savefolder = filedialog.askdirectory(initialdir = "S:\\Active Files\\PHOTOS\\STEWARDSHIP PHOTOS")
    destPath = savefolder
    outputPath.insert(0, destPath)

#The meat of the program. When you click "Save As" the actual PowerPoint file is
#initialized and the files are pulled from the source directory, sorted and placed
#into slides. The image is put into the image placeholder and the caption is pulled
#from the file name of the .jpg file. Then the PowerPoint is saved to the destination
#directory with a user supplied file name.
def saveClick():
    global sourcePath
    folderPath = sourcePath + '\\*.jpg'
    prs = Presentation()
    image_slide = 0
    filelist = glob.glob(folderPath)
    for filepath in natsorted(filelist, key = lambda fname: fname.lower()):
        pindex = filepath.rfind("\\")
        slide_layout = prs.slide_layouts[image_slide]
        slide = prs.slides.add_slide(slide_layout)
        placeholder = slide.placeholders[12]
        picture = placeholder.insert_picture(filepath)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            text_frame.text = (filepath[(pindex + 1):-4])
    savePath = destPath + "\\" + saveEntry.get() + ".pptx"
    prs.save(savePath)
    root.destroy()

#---------
#GUI Stuff
#---------

#initialize root
root = Tk()

#create folder path label
inputLabel = Label(root, text ='Select photo folder:')
inputLabel.grid(row = 2, column = 0, sticky = W)

#create text entry box for photo folder path
inputPath = Entry(root, width = 100)
inputPath.grid(row = 3, column = 0, padx = 10, pady = 20, columnspan = 2)

#create button to select source path
selectButton_1 = Button(root, text = "Select", width = 10, command = select_1Click)
selectButton_1.grid(row = 3, column = 2, padx = 10)

#create button to select save path
selectButton_2 = Button(root, text = "Select", width = 10, command = select_2Click)
selectButton_2.grid(row = 5, column = 2, padx = 10)

#create save path label
outputLabel = Label(root, text = 'Select folder to save presentation:')
outputLabel.grid(row = 4, column = 0, sticky = W)

#create text entry box for save path for presentation
outputPath = Entry(root, width = 100)
outputPath.grid(row = 5, column = 0, padx = 10, pady = 20, columnspan = 2)

#create save as label
saveName = Label(root, text = 'Save as:')
saveName.grid(row = 6, column = 0, sticky = W)

#create save as textbox
saveEntry = Entry(root, width = 100)
saveEntry.grid(row = 7, column = 0, columnspan = 2, padx = 10, pady = 20)

#create button to save presentation as
saveButton = Button(root, text = "Save As", width = 10, command = saveClick)
saveButton.grid(row = 7, column = 2, padx = 10)

#run mainloop
root.mainloop()

